#!/usr/bin/env python3
"""Generate broadcast-quality 1080p award certificate PNGs for Cub Scouts."""

import argparse
import csv
import json
import os
import sys
import textwrap
import urllib.request
from collections import OrderedDict
from PIL import Image, ImageDraw, ImageFont

# --- Configuration ---
WIDTH, HEIGHT = 1920, 1080
BG_COLOR = (26, 26, 46)  # dark navy
BAR_HEIGHT = 15
HEADER_AREA_TOP = BAR_HEIGHT + 20
HEADER_AREA_HEIGHT = 140
BODY_TOP = HEADER_AREA_TOP + HEADER_AREA_HEIGHT
BODY_BOTTOM = HEIGHT - BAR_HEIGHT - 20
BODY_HEIGHT = BODY_BOTTOM - BODY_TOP
SIDE_MARGIN = 80

RANK_COLORS = {
    "lions":   {"primary": (255, 199, 44),  "accent": (0, 63, 135),  "bg": (51, 40, 9),  "logo": "rank_lion.png"},
    "tigers":  {"primary": (252, 106, 33),  "accent": (0, 63, 135),  "bg": (50, 21, 7),  "logo": "rank_tiger.png"},
    "wolves":  {"primary": (188, 208, 54),  "accent": (0, 63, 135),  "bg": (30, 33, 9),  "logo": "rank_wolf.png"},
    "bears":   {"primary": (0, 174, 239),   "accent": (0, 63, 135),  "bg": (0, 28, 38),  "logo": "rank_bear.png"},
    "webelos": {"primary": (0, 132, 61),    "accent": (0, 63, 135),  "bg": (0, 26, 12),  "logo": "rank_webelos.png"},
    "aol":     {"primary": (0, 132, 61),    "accent": (0, 63, 135),  "bg": (0, 26, 12),  "logo": "rank_webelos.png"},
}

RANK_DISPLAY = {
    "lions": "Lion",
    "tigers": "Tiger",
    "wolves": "Wolf",
    "bears": "Bear",
    "webelos": "Webelos",
    "aol": "Arrow of Light",
}

RANK_ORDER = ["lions", "tigers", "wolves", "bears", "webelos", "aol"]

FONT_SEARCH_PATHS = [
    # msttcorefonts (Debian/Ubuntu: sudo apt install ttf-mscorefonts-installer)
    "/usr/share/fonts/truetype/msttcorefonts",
    # Common Linux paths
    "/usr/share/fonts/truetype",
    "/usr/share/fonts",
    "/usr/local/share/fonts",
    # macOS
    "/Library/Fonts",
    "/System/Library/Fonts",
    os.path.expanduser("~/Library/Fonts"),
    # Windows
    os.path.join(os.environ.get("WINDIR", "C:\\Windows"), "Fonts"),
]


def _find_font(names):
    """Search for a font file by trying multiple names across system font dirs."""
    for directory in FONT_SEARCH_PATHS:
        for name in names:
            path = os.path.join(directory, name)
            if os.path.exists(path):
                return path
    return None


FONT_IMPACT = _find_font([
    "Impact.ttf", "impact.ttf",
]) or _find_font([
    "Arial_Bold.ttf", "Arial Bold.ttf", "arialbd.ttf",
    "DejaVuSans-Bold.ttf",
])

FONT_BOLD = _find_font([
    "Arial_Bold.ttf", "Arial Bold.ttf", "arialbd.ttf",
    "DejaVuSans-Bold.ttf", "LiberationSans-Bold.ttf",
])

FONT_REGULAR = _find_font([
    "Arial.ttf", "arial.ttf",
    "DejaVuSans.ttf", "LiberationSans-Regular.ttf",
])

if not FONT_BOLD:
    print("Warning: No bold font found. Install ttf-mscorefonts-installer or dejavu fonts.")
    FONT_BOLD = ""
if not FONT_IMPACT:
    FONT_IMPACT = FONT_BOLD
if not FONT_REGULAR:
    FONT_REGULAR = FONT_BOLD

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMAGES_DIR = os.path.join(SCRIPT_DIR, "images")


def load_award_images(path):
    """Load award_images.json and build SKU -> {name, imageUrl400} dict."""
    with open(path) as f:
        data = json.load(f)
    sku_map = {}
    for item in data["items"]:
        sku_map[item["sku"]] = item
    return sku_map


def load_csv(path):
    """Parse CSV, return OrderedDict of (first, last) -> list of award dicts."""
    scouts = OrderedDict()
    with open(path, newline="") as f:
        reader = csv.DictReader(f)
        for row in reader:
            key = (row["First Name"].strip(), row["Last Name"].strip())
            if key not in scouts:
                scouts[key] = {
                    "first": key[0],
                    "last": key[1],
                    "den_type": row["Den Type"].strip(),
                    "den_number": row["Den Number"].strip(),
                    "awards": [],
                }
            scouts[key]["awards"].append({
                "sku": row["SKU"].strip(),
                "item_name": row["Item Name"].strip(),
                "item_type": row["Item Type"].strip(),
            })
    return scouts


def download_images(scouts, sku_map):
    """Download badge images for all unique SKUs, caching in images/ dir."""
    os.makedirs(IMAGES_DIR, exist_ok=True)
    all_skus = set()
    for scout in scouts.values():
        for award in scout["awards"]:
            all_skus.add(award["sku"])

    for sku in sorted(all_skus):
        dest = os.path.join(IMAGES_DIR, f"sku_{sku}.png")
        if os.path.exists(dest):
            continue

        source = sku_map.get(sku)
        if source:
            name = source.get("name", f"SKU {sku}")
            url = source.get("imageUrl400", "")
            try:
                print(f"  Downloading SKU {sku} ({name})...")
                urllib.request.urlretrieve(url, dest)
            except Exception as e:
                print(f"  Warning: Failed to download SKU {sku}: {e}")
                generate_placeholder(dest, name, sku)
        else:
            print(f"  No image source for SKU {sku}, generating placeholder...")
            generate_placeholder(dest, f"SKU {sku}", sku)


def generate_placeholder(path, name, sku):
    """Generate a simple circular placeholder badge."""
    size = 400
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Draw circle
    margin = 20
    draw.ellipse(
        [margin, margin, size - margin, size - margin],
        fill=(0, 63, 135),
        outline=(255, 199, 44),
        width=4,
    )
    # Draw text
    font = ImageFont.truetype(FONT_BOLD, 32)
    lines = name.split("\n") if "\n" in name else textwrap.wrap(name, width=14)
    y = size // 2 - len(lines) * 20
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        tw = bbox[2] - bbox[0]
        draw.text(((size - tw) // 2, y), line, fill="white", font=font)
        y += 40
    img.save(path)


def clean_award_name(name):
    """Clean up award names for display."""
    if name.endswith(" Adventure"):
        name = name[:-10]
    if name.endswith(" Emblem"):
        name = name[:-7]
    return name


def generate_scout_image(scout):
    """Generate a 1920x1080 PNG for one scout.

    Layout: each award is a horizontal row — small badge on left, large
    award name on right.  Awards are arranged in 1 or 2 columns depending
    on count.
    """
    den_type = scout["den_type"]
    colors = RANK_COLORS.get(den_type, RANK_COLORS["wolves"])
    primary = colors["primary"]
    accent = colors["accent"]
    bg = colors.get("bg", BG_COLOR)

    img = Image.new("RGB", (WIDTH, HEIGHT), bg)
    draw = ImageDraw.Draw(img)
    rank_name = RANK_DISPLAY.get(den_type, den_type.title())

    # --- Top and bottom color bars ---
    draw.rectangle([0, 0, WIDTH, BAR_HEIGHT], fill=primary)
    draw.rectangle([0, HEIGHT - BAR_HEIGHT, WIDTH, HEIGHT], fill=primary)

    # --- Scout name (left-aligned, Impact, uppercase) ---
    name_font = ImageFont.truetype(FONT_IMPACT, 90)
    full_name = f"{scout['first']} {scout['last']}".upper()
    name_x = SIDE_MARGIN
    name_y = HEADER_AREA_TOP + 5
    draw.text((name_x, name_y), full_name, fill="white", font=name_font)

    # --- Den info (left-aligned, below name with spacing) ---
    info_font = ImageFont.truetype(FONT_BOLD, 32)
    den_info = f"{rank_name}  \u2022  Den {scout['den_number']}"
    info_x = SIDE_MARGIN
    info_y = name_y + 100
    draw.text((info_x, info_y), den_info, fill=primary, font=info_font)

    # --- Split awards into featured (non-adventure) and adventures ---
    featured = [a for a in scout["awards"] if a["item_type"] != "Adventure"]
    adventures = [a for a in scout["awards"] if a["item_type"] == "Adventure"]

    body_top = BODY_TOP
    featured_height = 0

    # --- Featured section: rank emblems, recruiter strip, etc. ---
    if featured:
        feat_img_h = 120  # max image height
        feat_padding = 20
        featured_height = feat_img_h + feat_padding * 2 + 40  # image + label + padding
        feat_font = ImageFont.truetype(FONT_BOLD, 28)

        # Lay out featured items horizontally, centered
        feat_gap = 60  # gap between featured items
        # Measure total width needed
        feat_items = []
        for award in featured:
            sku = award["sku"]
            badge_path = os.path.join(IMAGES_DIR, f"sku_{sku}.png")
            display_name = clean_award_name(award["item_name"])
            bbox = draw.textbbox((0, 0), display_name, font=feat_font)
            label_w = bbox[2] - bbox[0]
            # Load image to get aspect ratio
            img_w = feat_img_h  # fallback square
            if os.path.exists(badge_path):
                try:
                    tmp = Image.open(badge_path)
                    scale = feat_img_h / tmp.height
                    img_w = int(tmp.width * scale)
                except Exception:
                    pass
            item_w = max(img_w, label_w)
            feat_items.append({"award": award, "img_w": img_w, "item_w": item_w, "label_w": label_w})

        total_feat_w = sum(f["item_w"] for f in feat_items) + feat_gap * (len(feat_items) - 1)
        feat_x = (WIDTH - total_feat_w) // 2
        feat_cy = body_top + feat_padding + feat_img_h // 2

        # Draw full-width colored band behind featured section
        band_color = tuple(min(255, c + 30) for c in bg)
        draw.rectangle([0, body_top, WIDTH, body_top + featured_height], fill=band_color)
        # Accent line at top and bottom of band
        draw.rectangle([0, body_top, WIDTH, body_top + 3], fill=primary)
        draw.rectangle([0, body_top + featured_height - 3, WIDTH, body_top + featured_height], fill=primary)

        for fi in feat_items:
            award = fi["award"]
            sku = award["sku"]
            badge_path = os.path.join(IMAGES_DIR, f"sku_{sku}.png")
            cx = feat_x + fi["item_w"] // 2

            # Draw image centered in item slot
            if os.path.exists(badge_path):
                try:
                    badge = Image.open(badge_path).convert("RGBA")
                    scale = feat_img_h / badge.height
                    bw = int(badge.width * scale)
                    badge = badge.resize((bw, feat_img_h), Image.LANCZOS)
                    img.paste(badge, (cx - bw // 2, feat_cy - feat_img_h // 2), badge)
                except Exception:
                    pass

            # Draw label centered below image
            display_name = clean_award_name(award["item_name"])
            bbox = draw.textbbox((0, 0), display_name, font=feat_font)
            lw = bbox[2] - bbox[0]
            label_y = feat_cy + feat_img_h // 2 + 8
            draw.text((cx - lw // 2, label_y), display_name, fill="white", font=feat_font)

            feat_x += fi["item_w"] + feat_gap

        body_top += featured_height

    # --- Adventure rows: image left, text right ---
    count = len(adventures)
    if count == 0:
        return img

    # Decide 1 or 2 columns
    if count <= 5:
        num_cols = 1
    else:
        num_cols = 2

    available_body_h = BODY_BOTTOM - body_top
    rows_per_col = (count + num_cols - 1) // num_cols
    col_width = (WIDTH - 2 * SIDE_MARGIN) // num_cols
    col_gap = 40

    row_height = available_body_h // rows_per_col
    badge_size = min(row_height - 12, 180)
    text_font_size = min(36, max(22, badge_size // 4 + 14))
    text_font = ImageFont.truetype(FONT_BOLD, text_font_size)
    text_gap = 20

    for idx, award in enumerate(adventures):
        col = idx // rows_per_col
        row = idx % rows_per_col

        if num_cols == 1:
            content_w = badge_size + text_gap + 500
            col_x = (WIDTH - content_w) // 2
        else:
            col_x = SIDE_MARGIN + col * (col_width + col_gap)

        # Both columns start at the same y, based on full left column height
        total_grid_h = rows_per_col * row_height
        grid_y_offset = body_top + (available_body_h - total_grid_h) // 2

        row_y = grid_y_offset + row * row_height
        row_cy = row_y + row_height // 2

        # Draw badge image (fit within badge_size box, preserve aspect ratio)
        sku = award["sku"]
        badge_path = os.path.join(IMAGES_DIR, f"sku_{sku}.png")
        if os.path.exists(badge_path):
            try:
                badge = Image.open(badge_path).convert("RGBA")
                badge.thumbnail((badge_size, badge_size), Image.LANCZOS)
                bw, bh = badge.size
                paste_x = col_x + (badge_size - bw) // 2
                paste_y = row_cy - bh // 2
                img.paste(badge, (paste_x, paste_y), badge)
            except Exception as e:
                print(f"  Warning: Could not load badge {badge_path}: {e}")

        # Draw award name to the right of badge area
        display_name = clean_award_name(award["item_name"])
        text_x = col_x + badge_size + text_gap
        bbox = draw.textbbox((0, 0), display_name, font=text_font)
        text_h = bbox[3] - bbox[1]
        text_y = row_cy - text_h // 2 - 2
        draw.text((text_x, text_y), display_name, fill="white", font=text_font)

    # --- Rank logo (right-aligned in header, drawn last so it's on top) ---
    logo_file = colors.get("logo")
    if logo_file:
        logo_path = os.path.join(IMAGES_DIR, logo_file)
        if os.path.exists(logo_path):
            logo = Image.open(logo_path).convert("RGBA")
            logo_max_h = HEADER_AREA_HEIGHT + 20
            scale = logo_max_h / logo.height
            logo_w = int(logo.width * scale)
            logo_h = logo_max_h
            logo = logo.resize((logo_w, logo_h), Image.LANCZOS)
            logo_x = WIDTH - SIDE_MARGIN - logo_w
            logo_y = HEADER_AREA_TOP
            img.paste(logo, (logo_x, logo_y), logo)

    return img


def generate_pptx(pptx_path, sorted_scouts):
    """Generate a PowerPoint with native shapes/text/images, one slide per scout."""
    from pptx import Presentation
    from pptx.util import Emu

    EMU_PX = 9525  # 1 pixel at 96 DPI

    prs = Presentation()
    prs.slide_width = Emu(WIDTH * EMU_PX)
    prs.slide_height = Emu(HEIGHT * EMU_PX)
    blank_layout = prs.slide_layouts[6]

    for scout in sorted_scouts:
        slide = prs.slides.add_slide(blank_layout)
        _render_scout_slide(slide, scout, EMU_PX)

    prs.save(pptx_path)


def _render_scout_slide(slide, scout, EMU_PX):
    """Render one scout's certificate as native PPTX shapes on a slide."""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    def px(val):
        return Emu(int(val * EMU_PX))

    def font_pt(pil_px_size):
        """Convert PIL pixel font size to PPTX points."""
        return Pt(pil_px_size * 72 / 96)

    def add_rect(x, y, w, h, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()

    def add_text(x, y, w, h, text, font_name, font_size, color,
                 bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
        tf = txBox.text_frame
        tf.auto_size = None
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.space_before = p.space_after = Pt(0)
        p.text = text
        p.font.size = font_pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(*color)
        p.font.bold = bold
        p.alignment = align

    def add_picture(path, x, y, max_w, max_h):
        """Add image preserving aspect ratio, centered in the given box."""
        if not os.path.exists(path):
            return
        try:
            with Image.open(path) as tmp:
                iw, ih = tmp.size
            scale = min(max_w / iw, max_h / ih)
            w, h = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(
                path,
                px(x + (max_w - w) // 2), px(y + (max_h - h) // 2),
                px(w), px(h),
            )
        except Exception:
            pass

    den_type = scout["den_type"]
    colors = RANK_COLORS.get(den_type, RANK_COLORS["wolves"])
    primary = colors["primary"]
    bg = colors.get("bg", BG_COLOR)
    rank_name = RANK_DISPLAY.get(den_type, den_type.title())

    # --- Slide background ---
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*bg)

    # --- Top and bottom color bars ---
    add_rect(0, 0, WIDTH, BAR_HEIGHT, primary)
    add_rect(0, HEIGHT - BAR_HEIGHT, WIDTH, BAR_HEIGHT, primary)

    # --- Scout name ---
    full_name = f"{scout['first']} {scout['last']}".upper()
    name_y = HEADER_AREA_TOP + 5
    add_text(SIDE_MARGIN, name_y, WIDTH - 2 * SIDE_MARGIN, 100,
             full_name, "Impact", 90, (255, 255, 255))

    # --- Den info ---
    den_info = f"{rank_name}  \u2022  Den {scout['den_number']}"
    info_y = name_y + 100
    add_text(SIDE_MARGIN, info_y, WIDTH - 2 * SIDE_MARGIN, 45,
             den_info, "Arial", 32, primary, bold=True)

    # --- Split awards ---
    featured = [a for a in scout["awards"] if a["item_type"] != "Adventure"]
    adventures = [a for a in scout["awards"] if a["item_type"] == "Adventure"]

    body_top = BODY_TOP

    # --- Featured section ---
    if featured:
        feat_img_h = 120
        feat_padding = 20
        featured_height = feat_img_h + feat_padding * 2 + 40

        # Use PIL for text/image measurement (same as PNG renderer)
        feat_font = ImageFont.truetype(FONT_BOLD, 28)
        _m = ImageDraw.Draw(Image.new("RGB", (1, 1)))

        feat_gap = 60
        feat_items = []
        for award in featured:
            sku = award["sku"]
            badge_path = os.path.join(IMAGES_DIR, f"sku_{sku}.png")
            display_name = clean_award_name(award["item_name"])
            bbox = _m.textbbox((0, 0), display_name, font=feat_font)
            label_w = bbox[2] - bbox[0]
            img_w = feat_img_h
            if os.path.exists(badge_path):
                try:
                    with Image.open(badge_path) as tmp:
                        img_w = int(tmp.width * feat_img_h / tmp.height)
                except Exception:
                    pass
            item_w = max(img_w, label_w)
            feat_items.append({"award": award, "img_w": img_w, "item_w": item_w})

        total_feat_w = sum(f["item_w"] for f in feat_items) + feat_gap * (len(feat_items) - 1)
        feat_x = (WIDTH - total_feat_w) // 2
        feat_cy = body_top + feat_padding + feat_img_h // 2

        # Band + accent lines
        band_color = tuple(min(255, c + 30) for c in bg)
        add_rect(0, body_top, WIDTH, featured_height, band_color)
        add_rect(0, body_top, WIDTH, 3, primary)
        add_rect(0, body_top + featured_height - 3, WIDTH, 3, primary)

        for fi in feat_items:
            award = fi["award"]
            badge_path = os.path.join(IMAGES_DIR, f"sku_{award['sku']}.png")

            add_picture(badge_path,
                        feat_x, feat_cy - feat_img_h // 2,
                        fi["item_w"], feat_img_h)

            display_name = clean_award_name(award["item_name"])
            label_y = feat_cy + feat_img_h // 2 + 8
            add_text(feat_x, label_y, fi["item_w"], 35,
                     display_name, "Arial", 28, (255, 255, 255),
                     bold=True, align=PP_ALIGN.CENTER)

            feat_x += fi["item_w"] + feat_gap

        body_top += featured_height

    # --- Adventure rows ---
    count = len(adventures)
    if count > 0:
        num_cols = 1 if count <= 5 else 2
        available_body_h = BODY_BOTTOM - body_top
        rows_per_col = (count + num_cols - 1) // num_cols
        col_width = (WIDTH - 2 * SIDE_MARGIN) // num_cols
        col_gap = 40
        row_height = available_body_h // rows_per_col
        badge_size = min(row_height - 12, 180)
        text_font_size = min(36, max(22, badge_size // 4 + 14))
        text_gap = 20

        for idx, award in enumerate(adventures):
            col = idx // rows_per_col
            row = idx % rows_per_col

            if num_cols == 1:
                content_w = badge_size + text_gap + 500
                col_x = (WIDTH - content_w) // 2
            else:
                col_x = SIDE_MARGIN + col * (col_width + col_gap)

            total_grid_h = rows_per_col * row_height
            grid_y_offset = body_top + (available_body_h - total_grid_h) // 2
            row_y = grid_y_offset + row * row_height
            row_cy = row_y + row_height // 2

            badge_path = os.path.join(IMAGES_DIR, f"sku_{award['sku']}.png")
            add_picture(badge_path,
                        col_x, row_cy - badge_size // 2,
                        badge_size, badge_size)

            display_name = clean_award_name(award["item_name"])
            text_x = col_x + badge_size + text_gap
            text_w = 500 if num_cols == 1 else col_width - badge_size - text_gap
            text_h = text_font_size + 10
            text_y = row_cy - text_h // 2
            add_text(text_x, text_y, text_w, text_h,
                     display_name, "Arial", text_font_size, (255, 255, 255),
                     bold=True)

    # --- Rank logo (last so it's on top) ---
    logo_file = colors.get("logo")
    if logo_file:
        logo_path = os.path.join(IMAGES_DIR, logo_file)
        if os.path.exists(logo_path):
            try:
                with Image.open(logo_path) as logo:
                    logo_max_h = HEADER_AREA_HEIGHT + 20
                    logo_w = int(logo.width * logo_max_h / logo.height)
                slide.shapes.add_picture(
                    logo_path,
                    px(WIDTH - SIDE_MARGIN - logo_w), px(HEADER_AREA_TOP),
                    px(logo_w), px(logo_max_h),
                )
            except Exception:
                pass


def parse_args():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Generate 1920x1080 award certificate PNGs for Cub Scouts.",
        epilog="The input CSV is exported from a Scoutbook purchase order. "
               "Required columns: First Name, Last Name, Den Type, Den Number, "
               "SKU, Item Type, Item Name.",
    )
    parser.add_argument(
        "csv", metavar="purchase_order.csv",
        help="CSV file exported from a Scoutbook purchase order",
    )
    parser.add_argument(
        "-o", "--output-dir", default=os.path.join(SCRIPT_DIR, "output"),
        help="directory to save generated PNGs (default: output/)",
    )
    parser.add_argument(
        "-p", "--pptx", metavar="FILE",
        help="also generate a PowerPoint presentation with all slides",
    )
    return parser.parse_args()


def main():
    args = parse_args()

    csv_path = args.csv
    if not os.path.isabs(csv_path):
        csv_path = os.path.join(os.getcwd(), csv_path)

    if not os.path.exists(csv_path):
        print(f"Error: File not found: {csv_path}")
        sys.exit(1)

    output_dir = args.output_dir
    if not os.path.isabs(output_dir):
        output_dir = os.path.join(os.getcwd(), output_dir)

    images_json_path = os.path.join(SCRIPT_DIR, "award_images.json")

    print("Loading award_images.json...")
    sku_map = load_award_images(images_json_path)
    print(f"  Loaded {len(sku_map)} award image entries")

    print(f"Loading CSV: {csv_path}")
    scouts = load_csv(csv_path)
    print(f"  Found {len(scouts)} scouts")

    print("Downloading/caching badge images...")
    download_images(scouts, sku_map)

    os.makedirs(output_dir, exist_ok=True)
    print(f"Generating certificate images in {output_dir}/...")

    # Sort scouts: by rank order, then last name within each rank
    rank_index = {r: i for i, r in enumerate(RANK_ORDER)}
    sorted_scouts = sorted(
        scouts.values(),
        key=lambda s: (rank_index.get(s["den_type"], 99), s["last"].lower(), s["first"].lower()),
    )

    for scout in sorted_scouts:
        img = generate_scout_image(scout)
        filename = f"{scout['den_type']}_{scout['last']}_{scout['first']}.png"
        out_path = os.path.join(output_dir, filename)
        img.save(out_path)
        award_count = len(scout["awards"])
        print(f"  {filename} ({award_count} award{'s' if award_count != 1 else ''})")

    print(f"\nDone! Generated {len(sorted_scouts)} images in {output_dir}/")

    if args.pptx:
        pptx_path = args.pptx
        if not os.path.isabs(pptx_path):
            pptx_path = os.path.join(os.getcwd(), pptx_path)
        print(f"Generating PowerPoint: {pptx_path}...")
        generate_pptx(pptx_path, sorted_scouts)
        print(f"  {len(sorted_scouts)} slides written.")


if __name__ == "__main__":
    main()
