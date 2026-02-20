"""Tests for scout_awards.renderer: PNG and PPTX generation."""

import os

import pytest
from PIL import Image

from scout_awards.data import load_csv, sort_scouts
from scout_awards.renderer import generate_pptx, generate_scout_image

TESTS_DIR = os.path.dirname(os.path.abspath(__file__))
GOLDEN_DIR = os.path.join(TESTS_DIR, "golden")


class TestGenerateScoutImage:
    def test_returns_image(self, sample_scout, images_dir):
        result = generate_scout_image(sample_scout, images_dir)
        assert isinstance(result, Image.Image)

    def test_dimensions(self, sample_scout, images_dir):
        result = generate_scout_image(sample_scout, images_dir)
        assert result.size == (1920, 1080)

    def test_mode(self, sample_scout, images_dir):
        result = generate_scout_image(sample_scout, images_dir)
        assert result.mode == "RGB"

    def test_all_ranks_render(self, images_dir):
        """Every rank produces a valid image without errors."""
        from scout_awards.config import RANK_ORDER
        for rank in RANK_ORDER:
            scout = {
                "first": "Test",
                "last": "Scout",
                "den_type": rank,
                "den_number": "1",
                "awards": [
                    {"sku": "619941", "item_name": "Test Adventure", "item_type": "Adventure"},
                ],
            }
            img = generate_scout_image(scout, images_dir)
            assert img.size == (1920, 1080)


@pytest.mark.slow
class TestGoldenFileRegression:
    """Compare generated images pixel-for-pixel against reference PNGs in test/.

    Run with: pytest -m slow
    Skip with: pytest -m 'not slow'
    """

    def _golden_files(self):
        """Return list of (scout, golden_path) tuples from sample.csv."""
        csv_path = os.path.join(TESTS_DIR, "sample.csv")
        scouts = sort_scouts(load_csv(csv_path))
        pairs = []
        for scout in scouts:
            filename = f"{scout['den_type']}_{scout['last']}_{scout['first']}.png"
            golden_path = os.path.join(GOLDEN_DIR, filename)
            if os.path.exists(golden_path):
                pairs.append((scout, golden_path))
        return pairs

    def test_golden_files_exist(self):
        pairs = self._golden_files()
        assert len(pairs) == 8, f"Expected 8 golden files, found {len(pairs)}"

    def test_pixel_identity(self, images_dir):
        pairs = self._golden_files()
        assert pairs, "No golden files found"
        for scout, golden_path in pairs:
            generated = generate_scout_image(scout, images_dir)
            golden = Image.open(golden_path).convert("RGB")
            name = f"{scout['first']} {scout['last']}"
            assert generated.size == golden.size, f"Size mismatch for {name}"
            assert generated.tobytes() == golden.tobytes(), f"Pixel mismatch for {name}"


class TestGeneratePptx:
    def test_creates_file(self, sample_scout, images_dir, tmp_path):
        pptx_path = str(tmp_path / "test.pptx")
        generate_pptx(pptx_path, [sample_scout], images_dir)
        assert os.path.exists(pptx_path)

    def test_slide_count(self, sample_csv_path, images_dir, tmp_path):
        scouts = sort_scouts(load_csv(sample_csv_path))
        pptx_path = str(tmp_path / "test.pptx")
        generate_pptx(pptx_path, scouts, images_dir)

        from pptx import Presentation
        prs = Presentation(pptx_path)
        assert len(prs.slides) == len(scouts)

    def test_slide_dimensions(self, sample_scout, images_dir, tmp_path):
        pptx_path = str(tmp_path / "test.pptx")
        generate_pptx(pptx_path, [sample_scout], images_dir)

        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation(pptx_path)
        EMU_PX = 9525
        assert prs.slide_width == Emu(1920 * EMU_PX)
        assert prs.slide_height == Emu(1080 * EMU_PX)
